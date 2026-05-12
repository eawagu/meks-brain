#!/usr/bin/env python3
"""
Convert any supported file to markdown. Output to stdout.

Usage: python converter.py <file_path>

Supported formats:
  Text:        .md, .txt, .csv, .tsv, .json, .xml, .yaml, .yml, .html, .htm, .log
  Documents:   .docx, .doc, .rtf, .epub
  Slides:      .pptx, .ppt
  Spreadsheet: .xlsx, .xls
  PDF:         .pdf
  Images:      .png, .jpg, .jpeg, .tiff, .bmp, .webp, .gif
  Email:       .eml, .msg
  Audio/Video: .mp3, .mp4, .m4a, .wav, .webm, .ogg, .flac
  Archives:    .zip, .tar.gz, .tar.bz2, .rar, .7z

Exit codes:
  0 = success (markdown on stdout)
  1 = error (message on stderr)
  2 = unknown format (extension on stderr)
"""

import sys
import os
import json
import subprocess
import tempfile
import email
from pathlib import Path

EXIT_SUCCESS = 0
EXIT_ERROR = 1
EXIT_UNKNOWN = 2


def passthrough(filepath: str) -> str:
    """Read text files as-is."""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def csv_to_md(filepath: str) -> str:
    """Convert CSV/TSV to markdown table."""
    import csv

    ext = Path(filepath).suffix.lower()
    delimiter = "\t" if ext in (".tsv",) else ","

    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f, delimiter=delimiter)
        rows = list(reader)

    if not rows:
        return ""

    header = rows[0]
    md = "| " + " | ".join(header) + " |\n"
    md += "| " + " | ".join(["---"] * len(header)) + " |\n"
    for row in rows[1:]:
        # Pad or truncate to header length
        padded = row + [""] * (len(header) - len(row))
        md += "| " + " | ".join(padded[: len(header)]) + " |\n"

    return md


def json_to_md(filepath: str) -> str:
    """Pretty-print JSON as fenced code block."""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        data = json.load(f)
    return f"```json\n{json.dumps(data, indent=2)}\n```"


def xml_to_md(filepath: str) -> str:
    """XML as fenced code block."""
    content = passthrough(filepath)
    return f"```xml\n{content}\n```"


def yaml_to_md(filepath: str) -> str:
    """YAML as fenced code block."""
    content = passthrough(filepath)
    return f"```yaml\n{content}\n```"


def html_to_md(filepath: str) -> str:
    """Convert HTML to markdown using pandoc if available, else strip tags."""
    pandoc_cmd = os.environ.get("PANDOC_CMD", "pandoc")
    try:
        result = subprocess.run(
            [pandoc_cmd, "-f", "html", "-t", "markdown", "--wrap=none", filepath],
            capture_output=True, timeout=60
        )
        if result.returncode == 0:
            return result.stdout.decode("utf-8", errors="replace")
    except FileNotFoundError:
        pass

    # Fallback: strip tags with basic regex
    import re
    content = passthrough(filepath)
    content = re.sub(r"<script[^>]*>.*?</script>", "", content, flags=re.DOTALL)
    content = re.sub(r"<style[^>]*>.*?</style>", "", content, flags=re.DOTALL)
    content = re.sub(r"<[^>]+>", "", content)
    return content.strip()


def _pandoc(fmt: str, filepath: str) -> str:
    """Convert a file to markdown via pandoc.

    No internal timeout — the outer wrapper in ingress.ts (10 min) bounds
    wall-clock time. Removing the inner 120s timeout eliminates size-based
    silent defer for docx files that pandoc could otherwise convert given
    a few more seconds. See designs/size-defer-fix.md.
    """
    pandoc_cmd = os.environ.get("PANDOC_CMD", "pandoc")
    result = subprocess.run(
        [pandoc_cmd, "-f", fmt, "-t", "markdown", "--wrap=none", filepath],
        capture_output=True
    )
    if result.returncode != 0:
        stderr = result.stderr.decode("utf-8", errors="replace") if result.stderr else ""
        raise RuntimeError(f"pandoc failed: {stderr}")
    return result.stdout.decode("utf-8", errors="replace")


def _docx_to_md_python_docx(filepath: str) -> str:
    """Lossy fallback: extract text content via python-docx when pandoc fails.

    No formatting fidelity (tables become pipe-joined rows, headings lose
    hierarchy, footnotes dropped) — but never fails on valid docx as long
    as the file's XML is well-formed. Used as the second-line path in
    docx_to_md so that pandoc errors never cause a silent defer to review/.

    Requires `python-docx` (pip install python-docx). If unavailable, the
    ImportError surfaces to the caller which converts it back into a
    conversion_failed → MISS calibration tuple — explicit signal, not silent.
    """
    from docx import Document
    doc = Document(filepath)
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def docx_to_md(filepath: str) -> str:
    """Convert DOCX to markdown.

    Primary path: pandoc (full formatting fidelity).
    Fallback path: python-docx (lossy text extraction) — used when pandoc
    errors so that no docx file is ever deferred to review/ purely on
    a pandoc-side failure mode.

    Either path's success means the file gets a source page. Both paths
    failing produces a RuntimeError that surfaces as conversion_failed in
    ingress.ts, which emits a MISS calibration tuple (no silent defer).
    """
    try:
        return _pandoc("docx", filepath)
    except RuntimeError as primary_err:
        try:
            return _docx_to_md_python_docx(filepath)
        except Exception as fallback_err:
            raise RuntimeError(
                f"pandoc failed: {primary_err}; "
                f"python-docx fallback also failed: {fallback_err}"
            )


def rtf_to_md(filepath: str) -> str:
    """Convert RTF to markdown via pandoc."""
    return _pandoc("rtf", filepath)


def epub_to_md(filepath: str) -> str:
    """Convert EPUB to markdown via pandoc."""
    return _pandoc("epub", filepath)


def legacy_office_to_md(filepath: str) -> str:
    """Convert legacy .doc/.ppt/.xls via LibreOffice headless → pandoc."""
    with tempfile.TemporaryDirectory() as tmpdir:
        # Convert to modern format via LibreOffice
        lo_cmd = os.environ.get("LIBREOFFICE_CMD", "libreoffice")
        result = subprocess.run(
            [lo_cmd, "--headless", "--convert-to", "docx", "--outdir", tmpdir, filepath],
            capture_output=True, timeout=120
        )
        if result.returncode != 0:
            stderr = result.stderr.decode("utf-8", errors="replace") if result.stderr else ""
            raise RuntimeError(f"LibreOffice conversion failed: {stderr}")

        # Find the converted file
        stem = Path(filepath).stem
        candidates = list(Path(tmpdir).glob(f"{stem}.*"))
        if not candidates:
            raise RuntimeError("LibreOffice produced no output file")

        converted = str(candidates[0])
        ext = candidates[0].suffix.lower()

        if ext == ".docx":
            return docx_to_md(converted)
        elif ext == ".pptx":
            return pptx_to_md(converted)
        elif ext == ".xlsx":
            return xlsx_to_md(converted)
        else:
            return passthrough(converted)


def pptx_to_md(filepath: str) -> str:
    """Convert PPTX to markdown using python-pptx."""
    from pptx import Presentation

    prs = Presentation(filepath)
    parts = []

    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {i}")

        # Slide title
        if slide.shapes.title and slide.shapes.title.text:
            parts.append(f"### {slide.shapes.title.text}")

        # All text shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and (not slide.shapes.title or text != slide.shapes.title.text):
                    parts.append(text)

            # Tables
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    md = "| " + " | ".join(rows[0]) + " |\n"
                    md += "| " + " | ".join(["---"] * len(rows[0])) + " |\n"
                    for row in rows[1:]:
                        padded = row + [""] * (len(rows[0]) - len(row))
                        md += "| " + " | ".join(padded[: len(rows[0])]) + " |\n"
                    parts.append(md)

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            notes = slide.notes_slide.notes_text_frame.text.strip()
            parts.append(f"> **Notes:** {notes}")

        parts.append("")

    return "\n\n".join(parts)


def xlsx_to_md(filepath: str) -> str:
    """Convert Excel to markdown tables (one per sheet)."""
    from openpyxl import load_workbook

    wb = load_workbook(filepath, read_only=True, data_only=True)
    parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        parts.append(f"## {sheet_name}")

        # Find actual data extent (skip fully empty rows)
        data_rows = [r for r in rows if any(c is not None for c in r)]
        if not data_rows:
            continue

        header = [str(c) if c is not None else "" for c in data_rows[0]]
        md = "| " + " | ".join(header) + " |\n"
        md += "| " + " | ".join(["---"] * len(header)) + " |\n"
        for row in data_rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            padded = cells + [""] * (len(header) - len(cells))
            md += "| " + " | ".join(padded[: len(header)]) + " |\n"

        parts.append(md)

    wb.close()
    return "\n\n".join(parts)


def pdf_to_md(filepath: str) -> str:
    """Convert PDF to markdown using pymupdf."""
    import fitz  # pymupdf

    doc = fitz.open(filepath)
    parts = []

    for i, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        if text:
            parts.append(f"<!-- Page {i} -->\n{text}")

    doc.close()
    return "\n\n".join(parts)


def image_to_md(filepath: str) -> str:
    """Images are handled natively by Claude vision via read_ingress — should not reach converter."""
    raise RuntimeError(
        f"Image files are processed via Claude vision, not converter.py. "
        f"This path should not be reached: {filepath}"
    )


def eml_to_md(filepath: str) -> str:
    """Parse .eml email to markdown."""
    with open(filepath, "rb") as f:
        msg = email.message_from_binary_file(f)

    parts = []
    parts.append(f"**From:** {msg.get('From', 'unknown')}")
    parts.append(f"**To:** {msg.get('To', 'unknown')}")
    parts.append(f"**Date:** {msg.get('Date', 'unknown')}")
    parts.append(f"**Subject:** {msg.get('Subject', 'unknown')}")
    parts.append("")

    # Body
    if msg.is_multipart():
        for part in msg.walk():
            content_type = part.get_content_type()
            if content_type == "text/plain":
                charset = part.get_content_charset() or "utf-8"
                payload = part.get_payload(decode=True)
                if payload:
                    parts.append(payload.decode(charset, errors="replace"))
            elif content_type == "text/html":
                charset = part.get_content_charset() or "utf-8"
                payload = part.get_payload(decode=True)
                if payload:
                    # Basic HTML strip
                    import re
                    html = payload.decode(charset, errors="replace")
                    text = re.sub(r"<[^>]+>", "", html)
                    parts.append(text.strip())
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            charset = msg.get_content_charset() or "utf-8"
            parts.append(payload.decode(charset, errors="replace"))

    # List attachments
    attachments = []
    if msg.is_multipart():
        for part in msg.walk():
            filename = part.get_filename()
            if filename:
                attachments.append(filename)
    if attachments:
        parts.append("\n**Attachments:** " + ", ".join(attachments))

    return "\n".join(parts)


def msg_to_md(filepath: str) -> str:
    """Parse Outlook .msg file to markdown."""
    import extract_msg

    msg = extract_msg.Message(filepath)
    parts = []
    parts.append(f"**From:** {msg.sender or 'unknown'}")
    parts.append(f"**To:** {msg.to or 'unknown'}")
    parts.append(f"**Date:** {msg.date or 'unknown'}")
    parts.append(f"**Subject:** {msg.subject or 'unknown'}")
    parts.append("")

    if msg.body:
        parts.append(msg.body)

    if msg.attachments:
        att_names = [a.longFilename or a.shortFilename or "unnamed" for a in msg.attachments]
        parts.append("\n**Attachments:** " + ", ".join(att_names))

    msg.close()
    return "\n".join(parts)


def audio_to_md(filepath: str) -> str:
    """Transcribe audio/video using OpenAI Whisper (local)."""
    # Try whisper CLI first
    whisper_cmd = os.environ.get("WHISPER_CMD", "whisper")
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                [whisper_cmd, filepath, "--output_format", "txt",
                 "--output_dir", tmpdir, "--language", "en"],
                capture_output=True, timeout=600
            )
            if result.returncode == 0:
                txt_files = list(Path(tmpdir).glob("*.txt"))
                if txt_files:
                    return txt_files[0].read_text(encoding="utf-8", errors="replace").strip()
    except FileNotFoundError:
        pass

    raise RuntimeError(
        "Whisper not available. Install with: pip install openai-whisper"
    )


def archive_to_md(filepath: str) -> str:
    """Extract archive and list contents with conversion attempt for each file."""
    import zipfile
    import tarfile

    ext = Path(filepath).suffix.lower()
    parts = [f"# Archive: {Path(filepath).name}\n"]

    with tempfile.TemporaryDirectory() as tmpdir:
        # Extract
        if ext == ".zip":
            with zipfile.ZipFile(filepath, "r") as zf:
                zf.extractall(tmpdir)
        elif ext in (".gz", ".bz2", ".xz") or filepath.endswith(
            (".tar.gz", ".tar.bz2", ".tar.xz", ".tgz")
        ):
            with tarfile.open(filepath, "r:*") as tf:
                tf.extractall(tmpdir, filter="data")
        else:
            # Try 7z CLI
            sevenz_cmd = os.environ.get("SEVENZ_CMD", "7z")
            result = subprocess.run(
                [sevenz_cmd, "x", filepath, f"-o{tmpdir}", "-y"],
                capture_output=True, timeout=120
            )
            if result.returncode != 0:
                stderr = result.stderr.decode("utf-8", errors="replace") if result.stderr else ""
                raise RuntimeError(f"Cannot extract archive: {stderr}")

        # Process each extracted file
        for root, _dirs, files in os.walk(tmpdir):
            for filename in sorted(files):
                fpath = os.path.join(root, filename)
                rel = os.path.relpath(fpath, tmpdir)
                parts.append(f"\n## {rel}\n")
                try:
                    content = convert(fpath)
                    parts.append(content)
                except SystemExit:
                    parts.append(f"*[Unsupported format: {Path(filename).suffix}]*")
                except Exception as e:
                    parts.append(f"*[Conversion error: {e}]*")

    return "\n".join(parts)


# ─── Format registry ──────────────────────────────────────────

CONVERTERS = {
    # Passthrough text
    ".md": passthrough,
    ".txt": passthrough,
    ".log": passthrough,
    # Structured text
    ".csv": csv_to_md,
    ".tsv": csv_to_md,
    ".json": json_to_md,
    ".xml": xml_to_md,
    ".yaml": yaml_to_md,
    ".yml": yaml_to_md,
    # HTML
    ".html": html_to_md,
    ".htm": html_to_md,
    # Documents
    ".docx": docx_to_md,
    ".rtf": rtf_to_md,
    ".epub": epub_to_md,
    ".doc": legacy_office_to_md,
    # Slides
    ".pptx": pptx_to_md,
    ".ppt": legacy_office_to_md,
    # Spreadsheets
    ".xlsx": xlsx_to_md,
    ".xls": legacy_office_to_md,
    # PDF
    ".pdf": pdf_to_md,
    # Images (OCR)
    ".png": image_to_md,
    ".jpg": image_to_md,
    ".jpeg": image_to_md,
    ".tiff": image_to_md,
    ".tif": image_to_md,
    ".bmp": image_to_md,
    ".webp": image_to_md,
    ".gif": image_to_md,
    # Email
    ".eml": eml_to_md,
    ".msg": msg_to_md,
    # Audio/Video — disabled, uncomment when whisper + ffmpeg are installed
    # ".mp3": audio_to_md,
    # ".mp4": audio_to_md,
    # ".m4a": audio_to_md,
    # ".wav": audio_to_md,
    # ".webm": audio_to_md,
    # ".ogg": audio_to_md,
    # ".flac": audio_to_md,
    # Archives
    ".zip": archive_to_md,
    ".rar": archive_to_md,
    ".7z": archive_to_md,
}

# Handle .tar.gz, .tar.bz2, etc. (double extensions)
DOUBLE_EXT_CONVERTERS = {
    ".tar.gz": archive_to_md,
    ".tar.bz2": archive_to_md,
    ".tar.xz": archive_to_md,
    ".tgz": archive_to_md,
}


def convert(filepath: str) -> str:
    """Convert a file to markdown. Returns the markdown string."""
    p = Path(filepath)

    if not p.exists():
        raise FileNotFoundError(f"File not found: {filepath}")

    # Check double extensions first
    name_lower = p.name.lower()
    for dext, converter in DOUBLE_EXT_CONVERTERS.items():
        if name_lower.endswith(dext):
            return converter(filepath)

    ext = p.suffix.lower()
    converter = CONVERTERS.get(ext)

    if converter is None:
        # Unknown format
        print(ext, file=sys.stderr)
        sys.exit(EXIT_UNKNOWN)

    return converter(filepath)


def main():
    if len(sys.argv) != 2:
        print("Usage: python converter.py <file_path>", file=sys.stderr)
        sys.exit(EXIT_ERROR)

    filepath = sys.argv[1]

    try:
        markdown = convert(filepath)
        print(markdown)
        sys.exit(EXIT_SUCCESS)
    except SystemExit:
        raise
    except Exception as e:
        print(f"Error converting {filepath}: {e}", file=sys.stderr)
        sys.exit(EXIT_ERROR)


if __name__ == "__main__":
    main()
